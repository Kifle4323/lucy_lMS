"""
PPTX to HTML Converter with Reading Time Tracker
Converts PowerPoint files to interactive HTML with slide-by-slide reading time tracking
"""

import os
import base64
import json
import sys
from io import BytesIO
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from datetime import datetime

# Fix Windows console encoding
if sys.platform == 'win32':
    import codecs
    sys.stdout = codecs.getwriter('utf-8')(sys.stdout.buffer)
    sys.stderr = codecs.getwriter('utf-8')(sys.stderr.buffer)

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


def calculate_reading_time_seconds(text):
    """
    Calculate required reading time based on text length
    Average reading speed: 200-250 words per minute = ~4 words per second
    """
    if not text:
        return 3  # Minimum 3 seconds for empty slides
    
    words = len(text.split())
    base_time = words * 0.3  # 0.3 seconds per word
    buffer_time = min(5, words * 0.1)  # Comprehension buffer
    total_time = base_time + buffer_time
    
    # Cap at reasonable limits
    if total_time < 5:
        total_time = 5
    elif total_time > 120:
        total_time = 120
    
    return round(total_time, 1)


def extract_slide_content(slide):
    """Extract text and images from a slide"""
    text_content = []
    images = []
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                
                format_type = 'png'
                if HAS_PIL:
                    try:
                        img = Image.open(BytesIO(image_bytes))
                        format_type = img.format.lower() if img.format else 'png'
                    except:
                        pass
                
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                images.append({
                    'format': format_type,
                    'data': image_base64
                })
            except Exception as e:
                print(f"Warning: Could not extract image: {e}", file=sys.stderr)
        
        elif shape.has_table:
            table_data = []
            for row in shape.table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            if table_data:
                text_content.append({'type': 'table', 'data': table_data})
        
        elif shape.has_text_frame:
            text = shape.text.strip()
            if text:
                is_title = False
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if shape.placeholder_format.type == 1:  # Title
                        is_title = True
                
                text_content.append({
                    'type': 'title' if is_title else 'text',
                    'content': text
                })
    
    full_text = ' '.join([item['content'] for item in text_content if item['type'] != 'table'])
    return text_content, images, full_text


def convert_pptx_to_html(pptx_file_path, material_id=""):
    """
    Convert PowerPoint to HTML with reading time tracking
    Returns the HTML content and slide metadata
    """
    
    if not os.path.exists(pptx_file_path):
        raise FileNotFoundError(f"File not found: {pptx_file_path}")
    
    prs = Presentation(pptx_file_path)
    total_slides = len(prs.slides)
    
    # Analyze each slide
    slides_data = []
    slides_html = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        text_content, images, full_text = extract_slide_content(slide)
        reading_time = calculate_reading_time_seconds(full_text)
        reading_time += len(images) * 2  # Add 2 seconds per image
        
        slide_data = {
            'slide_num': slide_num,
            'word_count': len(full_text.split()),
            'image_count': len(images),
            'required_time': reading_time,
            'text_preview': full_text[:100] + '...' if len(full_text) > 100 else full_text
        }
        slides_data.append(slide_data)
        
        # Build slide HTML
        content_html = []
        for item in text_content:
            if item['type'] == 'title':
                content_html.append(f"<h1>{item['content']}</h1>")
            elif item['type'] == 'table':
                table_rows = []
                for row in item['data']:
                    cells = ''.join([f"<td>{cell}</td>" for cell in row])
                    table_rows.append(f"<tr>{cells}</tr>")
                content_html.append(f"<table>{''.join(table_rows)}</table>")
            else:
                # Handle lists (lines with bullet points)
                content = item['content']
                if '\n' in content:
                    items = ''.join([f"<li>{line.strip()}</li>" for line in content.split('\n') if line.strip()])
                    content_html.append(f"<ul>{items}</ul>")
                else:
                    content_html.append(f"<p>{content}</p>")
        
        # Add images
        for img in images:
            content_html.append(
                f'<div class="image-container">'
                f'<img src="data:image/{img["format"]};base64,{img["data"]}" class="slide-image">'
                f'</div>'
            )
        
        slides_html.append({
            'slide_num': slide_num,
            'content': ''.join(content_html),
            'required_time': reading_time
        })
    
    return {
        'slides_data': slides_data,
        'slides_html': slides_html,
        'total_slides': total_slides,
        'title': Path(pptx_file_path).stem
    }


def generate_html_output(conversion_result, material_id, user_id="student"):
    """Generate the full HTML file with tracking"""
    
    slides_json = json.dumps(conversion_result['slides_data'])
    slides_html = conversion_result['slides_html']
    total_slides = conversion_result['total_slides']
    title = conversion_result['title']
    
    # Build slides HTML
    slides_content = []
    for slide in slides_html:
        slides_content.append(f'''
        <div class="slide unread" id="slide-{slide['slide_num']}" 
             data-slide-num="{slide['slide_num']}"
             data-required-time="{slide['required_time']}">
            <div class="slide-header">
                <span class="slide-number">Slide {slide['slide_num']} of {total_slides}</span>
            </div>
            <div class="slide-content">
                {slide['content']}
            </div>
        </div>
        ''')
    
    html = f'''<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        
        body {{
            font-family: 'Segoe UI', Arial, sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
            padding: 30px 20px;
        }}
        
        .header {{
            position: fixed;
            top: 0;
            left: 0;
            right: 0;
            background: white;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
            z-index: 1000;
            padding: 15px 30px;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }}
        
        .progress-container {{
            flex: 1;
            max-width: 600px;
            margin: 0 20px;
        }}
        
        .progress-bar {{
            width: 100%;
            height: 10px;
            background: #e0e0e0;
            border-radius: 5px;
            overflow: hidden;
        }}
        
        .progress-fill {{
            height: 100%;
            background: linear-gradient(90deg, #4CAF50, #8BC34A);
            width: 0%;
            transition: width 0.3s ease;
        }}
        
        .progress-text {{
            margin-top: 5px;
            font-size: 14px;
            color: #666;
        }}
        
        .slide-status {{
            display: flex;
            gap: 10px;
            align-items: center;
        }}
        
        .status-badge {{
            padding: 5px 15px;
            border-radius: 20px;
            font-size: 14px;
            font-weight: 600;
        }}
        
        .status-unread {{ background: #ff9800; color: white; }}
        .status-reading {{ background: #2196F3; color: white; }}
        .status-read {{ background: #4CAF50; color: white; }}
        
        .timer-display {{
            font-family: monospace;
            font-size: 18px;
            padding: 5px 15px;
            background: #f5f5f5;
            border-radius: 20px;
        }}
        
        .presentation {{
            max-width: 1200px;
            margin: 80px auto 30px;
        }}
        
        .slide {{
            background: white;
            margin: 30px auto;
            padding: 40px;
            border-radius: 15px;
            box-shadow: 0 10px 30px rgba(0,0,0,0.2);
            position: relative;
            transition: border 0.3s ease;
        }}
        
        .slide.unread {{ border-left: 5px solid #ff9800; }}
        .slide.reading {{ border-left: 5px solid #2196F3; }}
        .slide.read {{ border-left: 5px solid #4CAF50; }}
        
        .slide-header {{
            display: flex;
            justify-content: space-between;
            align-items: center;
            margin-bottom: 15px;
            padding-bottom: 8px;
            border-bottom: 1px solid #eee;
        }}
        
        .slide-number {{
            font-size: 16px;
            color: #999;
        }}
        
        
        .slide-content {{ margin-top: 20px; }}
        
        h1 {{
            color: #2c3e50;
            font-size: 2.5em;
            margin-bottom: 20px;
        }}
        
        h2 {{
            color: #34495e;
            font-size: 2em;
            margin: 15px 0;
        }}
        
        p {{
            font-size: 1.1em;
            line-height: 1.6;
            color: #555;
            margin: 15px 0;
        }}
        
        ul, ol {{
            font-size: 1.1em;
            line-height: 1.8;
            margin: 15px 0 15px 30px;
        }}
        
        li {{ margin: 8px 0; }}
        
        .slide-image {{
            max-width: 100%;
            height: auto;
            margin: 20px 0;
            border-radius: 8px;
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
        }}
        
        .image-container {{
            text-align: center;
            margin: 20px 0;
        }}
        
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }}
        
        td, th {{
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }}
        
        .navigation {{
            display: flex;
            justify-content: center;
            gap: 15px;
            margin: 30px 0;
            flex-wrap: wrap;
        }}
        
        .nav-btn {{
            background: white;
            border: none;
            padding: 12px 24px;
            border-radius: 30px;
            cursor: pointer;
            font-size: 16px;
            font-weight: 600;
            color: #667eea;
            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
            transition: all 0.3s ease;
        }}
        
        .nav-btn:hover {{
            background: #667eea;
            color: white;
            transform: translateY(-2px);
        }}
        
        .nav-btn:disabled {{
            opacity: 0.5;
            cursor: not-allowed;
        }}
        
        .tracking-panel {{
            position: fixed;
            bottom: 20px;
            right: 20px;
            background: white;
            padding: 20px;
            border-radius: 15px;
            box-shadow: 0 5px 20px rgba(0,0,0,0.2);
            max-width: 300px;
        }}
        
        .tracking-title {{
            font-weight: bold;
            margin-bottom: 10px;
            color: #333;
        }}
        
        .slide-progress-item {{
            display: flex;
            align-items: center;
            margin: 8px 0;
            font-size: 13px;
        }}
        
        .progress-indicator {{
            width: 12px;
            height: 12px;
            border-radius: 50%;
            margin-right: 8px;
        }}
        
        .indicator-unread {{ background: #ff9800; }}
        .indicator-reading {{ background: #2196F3; }}
        .indicator-read {{ background: #4CAF50; }}
        
        @media (max-width: 768px) {{
            .header {{ flex-direction: column; padding: 10px; }}
            .progress-container {{ margin: 10px 0; }}
            .slide {{ padding: 20px; }}
        }}
    </style>
</head>
<body>
    <div class="header">
        <h3>{title}</h3>
        <div class="progress-container">
            <div class="progress-bar">
                <div class="progress-fill" id="overallProgress"></div>
            </div>
            <div class="progress-text" id="progressText">0% Complete</div>
        </div>
        <div class="slide-status">
            <div class="timer-display" id="globalTimer">00:00</div>
            <div class="status-badge status-unread" id="globalStatus">Not Started</div>
        </div>
    </div>
    
    <div class="presentation" id="presentation">
        {''.join(slides_content)}
    </div>
    
    
    <div class="navigation">
        <button class="nav-btn" onclick="scrollToSlide(1)">First</button>
        <button class="nav-btn" onclick="previousSlide()">Previous</button>
        <button class="nav-btn" onclick="nextSlide()">Next</button>
        <button class="nav-btn" onclick="scrollToSlide({total_slides})">Last</button>
    </div>
    
    <script>
        const slidesData = {slides_json};
        const totalSlides = {total_slides};
        const materialId = "{material_id}";
        const apiBase = window.parent.API_BASE || '';
        
        let currentSlide = 1;
        let globalTimer = 0;
        let timerInterval = null;
        let slideTimeSpent = {{}};
        let slideStatuses = {{}};
        
        slidesData.forEach(slide => {{
            slideStatuses[slide.slide_num] = 'unread';
            slideTimeSpent[slide.slide_num] = 0;
        }});
        
        function startGlobalTimer() {{
            if (timerInterval) return;
            timerInterval = setInterval(() => {{
                globalTimer++;
                updateGlobalTimerDisplay();
                if (currentSlide && slideStatuses[currentSlide] !== 'read') {{
                    addTimeToSlide(currentSlide, 1);
                }}
            }}, 1000);
        }}
        
        function updateGlobalTimerDisplay() {{
            const minutes = Math.floor(globalTimer / 60);
            const seconds = globalTimer % 60;
            document.getElementById('globalTimer').textContent = 
                `${{minutes.toString().padStart(2, '0')}}:${{seconds.toString().padStart(2, '0')}}`;
        }}
        
        function addTimeToSlide(slideNum, seconds) {{
            if (slideStatuses[slideNum] === 'read') return;
            
            const slideData = slidesData.find(s => s.slide_num === slideNum);
            if (!slideData) return;
            
            slideTimeSpent[slideNum] = (slideTimeSpent[slideNum] || 0) + seconds;
            
            if (slideTimeSpent[slideNum] >= 1 && slideStatuses[slideNum] === 'unread') {{
                updateSlideStatus(slideNum, 'reading');
            }}
            
            if (slideTimeSpent[slideNum] >= slideData.required_time && slideStatuses[slideNum] !== 'read') {{
                updateSlideStatus(slideNum, 'read');
            }}
            
            updateProgressDisplay();
        }}
        
        function updateSlideStatus(slideNum, status) {{
            slideStatuses[slideNum] = status;
            
            const slideElement = document.getElementById(`slide-${{slideNum}}`);
            if (slideElement) {{
                slideElement.classList.remove('unread', 'reading', 'read');
                slideElement.classList.add(status);
            }}
            
            updateGlobalStatus();
            
            // Auto-save when slide is completed
            if (status === 'read') {{
                autoSaveProgress();
            }}
        }}
        
        let autoSaveTimeout = null;
        function autoSaveProgress() {{
            // Debounce auto-save to avoid too many requests
            if (autoSaveTimeout) clearTimeout(autoSaveTimeout);
            autoSaveTimeout = setTimeout(() => {{
                saveProgress(true); // silent save
            }}, 2000);
        }}
        
        function updateGlobalStatus() {{
            const readCount = Object.values(slideStatuses).filter(s => s === 'read').length;
            const readingCount = Object.values(slideStatuses).filter(s => s === 'reading').length;
            
            const statusElement = document.getElementById('globalStatus');
            const progressFill = document.getElementById('overallProgress');
            const progressText = document.getElementById('progressText');
            
            const percentage = Math.round((readCount / totalSlides) * 100);
            progressFill.style.width = `${{percentage}}%`;
            progressText.textContent = `${{percentage}}% Complete (${{readCount}}/${{totalSlides}} slides read)`;
            
            if (readCount === totalSlides) {{
                statusElement.textContent = 'Completed';
                statusElement.className = 'status-badge status-read';
            }} else if (readingCount > 0) {{
                statusElement.textContent = 'Reading';
                statusElement.className = 'status-badge status-reading';
            }} else {{
                statusElement.textContent = 'Not Started';
                statusElement.className = 'status-badge status-unread';
            }}
        }}
        
        function updateProgressDisplay() {{
            updateGlobalStatus();
        }}
        
        async function saveProgress(silent = false) {{
            const progress = {{
                materialId: materialId,
                totalTime: globalTimer,
                slideStatuses: slideStatuses,
                slideTimeSpent: slideTimeSpent,
                completedSlides: Object.values(slideStatuses).filter(s => s === 'read').length,
                totalSlides: totalSlides,
                timestamp: new Date().toISOString()
            }};
            
            try {{
                const response = await fetch(`${{apiBase}}/materials/${{materialId}}/progress`, {{
                    method: 'POST',
                    headers: {{ 'Content-Type': 'application/json' }},
                    body: JSON.stringify(progress)
                }});
                
                if (!response.ok) {{
                    throw new Error('Failed to save');
                }}
                
                if (!silent) {{
                    console.log('Progress saved');
                }}
            }} catch (err) {{
                // Fallback to localStorage
                localStorage.setItem(`ppt_progress_${{materialId}}`, JSON.stringify(progress));
                if (!silent) {{
                    console.log('Progress saved locally');
                }}
            }}
        }}
        
        async function loadProgress() {{
            try {{
                const response = await fetch(`${{apiBase}}/materials/${{materialId}}/progress`);
                if (response.ok) {{
                    const progress = await response.json();
                    restoreProgress(progress);
                    return;
                }}
            }} catch (err) {{
                console.log('Server load failed, trying localStorage');
            }}
            
            // Try localStorage
            const saved = localStorage.getItem(`ppt_progress_${{materialId}}`);
            if (saved) {{
                try {{
                    restoreProgress(JSON.parse(saved));
                }} catch (e) {{}}
            }}
        }}
        
        function restoreProgress(progress) {{
            if (!progress) return;
            
            if (progress.slideTimeSpent) {{
                Object.entries(progress.slideTimeSpent).forEach(([slideNum, time]) => {{
                    slideTimeSpent[parseInt(slideNum)] = time;
                }});
            }}
            
            if (progress.slideStatuses) {{
                Object.entries(progress.slideStatuses).forEach(([slideNum, status]) => {{
                    updateSlideStatus(parseInt(slideNum), status);
                }});
            }}
            
            globalTimer = progress.totalTime || 0;
            updateGlobalTimerDisplay();
            updateProgressDisplay();
        }}
        
        function scrollToSlide(slideNum) {{
            const slide = document.getElementById(`slide-${{slideNum}}`);
            if (slide) {{
                slide.scrollIntoView({{ behavior: 'smooth', block: 'start' }});
                currentSlide = slideNum;
            }}
        }}
        
        function nextSlide() {{
            if (currentSlide < totalSlides) {{
                scrollToSlide(currentSlide + 1);
            }}
        }}
        
        function previousSlide() {{
            if (currentSlide > 1) {{
                scrollToSlide(currentSlide - 1);
            }}
        }}
        
        const observer = new IntersectionObserver((entries) => {{
            entries.forEach(entry => {{
                if (entry.isIntersecting) {{
                    const slideId = entry.target.id;
                    const slideNum = parseInt(slideId.split('-')[1]);
                    currentSlide = slideNum;
                }}
            }});
        }}, {{ threshold: 0.5 }});
        
        document.addEventListener('DOMContentLoaded', () => {{
            document.querySelectorAll('.slide').forEach(slide => {{
                observer.observe(slide);
            }});
            
            startGlobalTimer();
            loadProgress();
            
            // Auto-save progress every 30 seconds
            setInterval(() => saveProgress(true), 30000);
            
            document.addEventListener('keydown', (e) => {{
                if (e.key === 'ArrowLeft') {{
                    e.preventDefault();
                    previousSlide();
                }} else if (e.key === 'ArrowRight' || e.key === ' ') {{
                    e.preventDefault();
                    nextSlide();
                }}
            }});
        }});
    </script>
</body>
</html>'''
    
    return html


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description='Convert PPTX to HTML with reading tracking')
    parser.add_argument('input', help='Input PPTX file path')
    parser.add_argument('output', help='Output HTML file path')
    parser.add_argument('--material-id', default='', help='Material ID for tracking')
    
    args = parser.parse_args()
    
    try:
        result = convert_pptx_to_html(args.input, args.material_id)
        html = generate_html_output(result, args.material_id)
        
        with open(args.output, 'w', encoding='utf-8') as f:
            f.write(html)
        
        print(f"Converted {args.input} to {args.output}")
        print(f"{result['total_slides']} slides, {sum(s['required_time'] for s in result['slides_data'])/60:.1f} min reading time")
        
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)
